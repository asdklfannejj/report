# report.py — 채팅형 보고서 메이커 / Plotly / Ollama 설명 / 추천 플랜(사람 친화형+편집) / PPT·Word·Excel
# 요구 패키지(requirements.txt):
# streamlit, pandas, numpy, plotly, kaleido, pdfplumber, openpyxl, xlsxwriter, python-pptx, python-docx, requests
import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.figure_factory as ff
import pdfplumber
import io, zipfile, requests, os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt
from pptx.enum.text import PP_ALIGN

from docx import Document
from docx.shared import Inches as DocxInches

from matplotlib import font_manager, rcParams
import matplotlib.pyplot as plt

# ---------- 한글 폰트 ----------
FONT_PATH = os.path.join(os.path.dirname(__file__), "fonts", "MaruBuri-Regular.ttf")
if os.path.exists(FONT_PATH):
    try:
        font_manager.fontManager.addfont(FONT_PATH)
        rcParams["font.family"] = "MaruBuri"
    except Exception:
        rcParams["font.family"] = ["Malgun Gothic", "AppleGothic", "Noto Sans CJK KR", "NanumGothic", "DejaVu Sans"]
else:
    rcParams["font.family"] = ["Malgun Gothic", "AppleGothic", "Noto Sans CJK KR", "NanumGothic", "DejaVu Sans"]
rcParams["axes.unicode_minus"] = False
PLOTLY_FONT_FAMILY = "MaruBuri, NanumGothic, 'Malgun Gothic', AppleGothic, 'Noto Sans CJK KR', 'DejaVu Sans', sans-serif"

# ---------- Plotly 스타일 ----------
def _style_plotly(fig, title=None):
    fig.update_layout(
        template="plotly_white",
        title=title if title else fig.layout.title.text,
        title_x=0.5,
        title_font_size=18,
        font=dict(family=PLOTLY_FONT_FAMILY),
        margin=dict(l=50, r=30, t=90, b=50),
    )
    fig.update_xaxes(title_standoff=8, automargin=True)
    fig.update_yaxes(title_standoff=12, automargin=True)
    return fig

def _fig_to_png_bytes(fig):
    try:
        return fig.to_image(format="png", scale=2)  # kaleido 필요
    except Exception:
        if not st.session_state.get("_warn_kaleido", False):
            st.session_state["_warn_kaleido"] = True
            st.info("PNG 내보내기용으로 `kaleido`가 필요합니다. 설치가 없으면 Matplotlib로 대체합니다.")
        return None

# ---------- KDE(선택; scipy 있으면 사용) ----------
def _try_kde(values, n_points=200):
    try:
        from numpy import isfinite, linspace
        from scipy.stats import gaussian_kde
        vals = np.asarray(values, dtype=float)
        vals = vals[isfinite(vals)]
        if vals.size < 2:
            return None, None
        kde = gaussian_kde(vals)
        xs = linspace(float(vals.min()), float(vals.max()), n_points)
        ys = kde(xs)
        return xs, ys
    except Exception:
        return None, None

# ---------- Matplotlib 대체 PNG ----------
def _mpl_hist_png(series, title, xlabel, bins=20):
    buf = io.BytesIO()
    fig, ax = plt.subplots(figsize=(6.4, 4.2))
    ax.hist(series.dropna().values, bins=bins, edgecolor="white", alpha=0.9)
    ax.set_title(title); ax.set_xlabel(xlabel); ax.set_ylabel("개수")
    fig.tight_layout(); fig.savefig(buf, format="png", dpi=200); plt.close(fig)
    return buf.getvalue()

def _mpl_corr_png(corr_df, title):
    buf = io.BytesIO()
    fig, ax = plt.subplots(figsize=(6.4, 5.0))
    cax = ax.imshow(corr_df.values, cmap="RdBu", vmin=-1, vmax=1)
    ax.set_xticks(np.arange(len(corr_df.columns))); ax.set_yticks(np.arange(len(corr_df.columns)))
    ax.set_xticklabels(corr_df.columns, rotation=90); ax.set_yticklabels(corr_df.columns)
    ax.set_title(title); fig.colorbar(cax); fig.tight_layout(); fig.savefig(buf, format="png", dpi=200); plt.close(fig)
    return buf.getvalue()

# ---------- Ollama ----------
@st.cache_data(show_spinner=False, ttl=600)
def _ollama_generate(base_url: str, model: str, prompt: str, timeout: int = 90) -> str:
    try:
        url = base_url.rstrip("/") + "/api/generate"
        payload = {"model": model, "prompt": prompt, "stream": False}
        r = requests.post(url, json=payload, timeout=timeout)
        r.raise_for_status()
        data = r.json()
        return data.get("response") or data.get("text") or ""
    except Exception as e:
        return f"(설명 생성 실패: {e})"

# ---------- 추천 플랜 ----------
def infer_schema(df: pd.DataFrame):
    date_cols, numeric_cols, cat_cols = [], [], []
    for c in df.columns:
        s = df[c]
        if pd.api.types.is_numeric_dtype(s):
            numeric_cols.append(c)
        else:
            try:
                pd.to_datetime(s, errors="raise", infer_datetime_format=True)
                date_cols.append(c)
            except Exception:
                if s.dropna().nunique() <= max(20, int(len(s) * 0.05)):
                    cat_cols.append(c)
    return date_cols, numeric_cols, cat_cols

def recommend_plan(df: pd.DataFrame) -> dict:
    date_cols, numeric_cols, cat_cols = infer_schema(df)
    plan = {"timeseries": [], "numeric_dists": [], "categoricals": [], "correlation": False}
    if date_cols and numeric_cols:
        for n in numeric_cols[:3]:
            plan["timeseries"].append((date_cols[0], n))
    plan["numeric_dists"] = numeric_cols[:5]
    plan["categoricals"] = cat_cols[:3]
    plan["correlation"] = len(numeric_cols) >= 2
    return plan

# --- 사람이 읽기 쉬운 추천 플랜 렌더러 & 편집기 ---
def _humanize_plan(plan: dict) -> str:
    ts = plan.get("timeseries", [])
    nd = plan.get("numeric_dists", [])
    ct = plan.get("categoricals", [])
    cor = plan.get("correlation", False)

    md = []
    md.append("### 🔎 추천 보고서 플랜 (요약)")
    if ts:
        md.append(f"- ⏱ **시계열**: 날짜 기준 지표 {len(ts)}개")
    if nd:
        md.append(f"- 📈 **숫자 분포**: {len(nd)}개 지표")
    if ct:
        md.append(f"- 🧩 **범주 분포**: {len(ct)}개 컬럼")
    md.append(f"- 🔗 **상관관계 분석**: {'실행 권장' if cor else '불필요'}")
    md.append("")

    if ts:
        md.append("**시계열 후보**")
        md.append("| 날짜열 | 지표 |")
        md.append("|---|---|")
        for dcol, ncol in ts:
            md.append(f"| {dcol} | {ncol} |")
        md.append("")

    if nd:
        md.append("**숫자형 분포 후보**")
        md.append(", ".join([f"`{c}`" for c in nd]))
        md.append("")

    if ct:
        md.append("**범주형 분포 후보**")
        md.append(", ".join([f"`{c}`" for c in ct]))
        md.append("")

    if not (ts or nd or ct):
        md.append("> 🤔 유의미한 추천이 없습니다. 최소 1개 이상의 숫자형/날짜형 컬럼이 필요해요.")
    return "\n".join(md)

def _plan_editor(plan: dict, key_prefix: str):
    """
    추천 플랜을 사용자가 바로 다듬을 수 있는 간단 편집기(체크/선택).
    key_prefix를 붙여서 위젯 키 중복 방지.
    """
    with st.expander("🛠 추천 플랜 편집", expanded=False):
        # 시계열 체크박스
        ts = plan.get("timeseries", [])
        if ts:
            st.markdown("**⏱ 시계열 포함 여부**")
            new_ts = []
            for (dcol, ncol) in ts:
                on = st.checkbox(
                    f"{dcol} → {ncol}",
                    value=True,
                    key=f"{key_prefix}_ts_{dcol}_{ncol}"
                )
                if on:
                    new_ts.append((dcol, ncol))
            plan["timeseries"] = new_ts

        # 숫자형 멀티셀렉트
        nd = plan.get("numeric_dists", [])
        if nd:
            st.markdown("**📈 숫자형 분포(멀티 선택)**")
            selected_nd = st.multiselect(
                "포함할 지표",
                options=nd,
                default=nd,
                key=f"{key_prefix}_nd_select"
            )
            plan["numeric_dists"] = selected_nd

        # 범주형 멀티셀렉트
        ct = plan.get("categoricals", [])
        if ct:
            st.markdown("**🧩 범주형 분포(멀티 선택)**")
            selected_ct = st.multiselect(
                "포함할 컬럼",
                options=ct,
                default=ct,
                key=f"{key_prefix}_ct_select"
            )
            plan["categoricals"] = selected_ct

        # 상관관계 체크박스
        plan["correlation"] = st.checkbox(
            "🔗 상관관계 분석 포함",
            value=plan.get("correlation", False),
            key=f"{key_prefix}_cor_on"
        )

        st.info("설정이 즉시 반영됩니다. ‘결과 보고서 생성’ 시 편집된 플랜이 사용됩니다.")
    return plan

# ---------- UI: 사이드바 ----------
st.title("✨ 채팅형 이벤트 결과보고서 메이커 (Ollama 설명 포함)")

with st.sidebar:
    st.header("📡 Ollama 설정")
    use_ollama = st.checkbox("그래프/요약 자동 설명 생성", value=False)
    ollama_model = st.text_input("모델", value="llama3.1")
    ollama_base = st.text_input("서버 URL", value="http://127.0.0.1:11500")
    st.caption("로컬 또는 원격 Ollama 서버를 사용합니다. (예: llama3.1, qwen2.5:7b-instruct)")

    st.header("📊 차트 옵션")
    BIN_MODE = st.radio("빈 구분", ["자동", "개수 지정", "간격 지정"], index=0, horizontal=True)
    nbins = st.slider("빈 개수", 5, 100, 20) if BIN_MODE == "개수 지정" else None
    binsize = st.number_input("빈 간격", min_value=0.0, value=0.0, step=1.0) if BIN_MODE == "간격 지정" else None
    bargap = st.slider("막대 간격", 0.00, 0.50, 0.25, 0.01)
    show_kde = st.checkbox("밀도 곡선(KDE)", value=True)
    y_scale = st.selectbox("세로축", ["count", "percent", "probability density"], index=0)

# ---------- 업로드 ----------
uploaded_files = st.file_uploader("📂 파일 업로드 (Excel 또는 PDF)", type=["xlsx","xls","pdf"], accept_multiple_files=True)

# ---------- 채팅 상태 ----------
if "messages" not in st.session_state:
    st.session_state.messages = []
if "prefs" not in st.session_state:
    st.session_state.prefs = {
        "audience": "일반",
        "tone": "간결",
        "outputs": {"md": True, "excel": True, "ppt": True, "docx": True},
        "sections": {"overview": True, "numeric": True, "categorical": True, "timeseries": True, "correlation": True},
        "kpis": [],
    }

def chat_bot_reply(user_text):
    u = user_text.strip().lower()
    prefs = st.session_state.prefs
    msg = ""
    if u.startswith("/audience"):
        if "임원" in user_text: prefs["audience"] = "임원"; msg = "대상: 임원용으로 설정했어요."
        elif "실무" in user_text: prefs["audience"] = "실무"; msg = "대상: 실무용으로 설정했어요."
        else: msg = "대상을 '일반/임원/실무' 중에서 말씀해 주세요."
    elif u.startswith("/tone"):
        if "상세" in user_text: prefs["tone"] = "상세"; msg = "어조: 상세 설명으로 설정했어요."
        else: prefs["tone"] = "간결"; msg = "어조: 간결한 요약으로 설정했어요."
    elif u.startswith("/out"):
        prefs["outputs"] = {"md": "md" in u or "markdown" in u,
                            "excel": "excel" in u or "xlsx" in u,
                            "ppt": "ppt" in u or "파워포인트" in u or "슬라이드" in u,
                            "docx": "word" in u or "docx" in u or "워드" in u}
        msg = f"산출물 설정: {prefs['outputs']}"
    elif u.startswith("/kpi"):
        cols = [c.strip() for c in user_text.split(" ",1)[1].split(",")] if " " in user_text else []
        prefs["kpis"] = [c for c in cols if c]
        msg = f"KPI 컬럼 지정: {prefs['kpis']}"
    else:
        if "임원" in u: prefs["audience"]="임원"; msg += "임원용 요약 위주로 구성할게요. "
        if "실무" in u: prefs["audience"]="실무"; msg += "실무용 상세지표 중심으로 구성할게요. "
        if "간결" in u: prefs["tone"]="간결"; msg += "간결한 서술로 정리합니다. "
        if "상세" in u: prefs["tone"]="상세"; msg += "상세한 서술로 정리합니다. "
        if "ppt" in u or "파워포인트" in u or "슬라이드" in u: prefs["outputs"]["ppt"] = True
        if "word" in u or "docx" in u or "워드" in u: prefs["outputs"]["docx"] = True
        if "excel" in u or "엑셀" in u or "xlsx" in u: prefs["outputs"]["excel"] = True
        if "markdown" in u or "md" in u: prefs["outputs"]["md"] = True
        msg += "설정 반영 완료! '/audience 임원', '/tone 상세', '/out ppt,word', '/kpi 매출,전환율'처럼도 지시할 수 있어요."
    return msg

# ---------- 유틸 ----------
def _y_axis_label_and_format(histnorm):
    if histnorm in (None, "count"): return "개수", ",.0f", None, "none"
    if histnorm == "percent": return "백분율(%)", ".1f", "%", "none"
    return "확률밀도", ".3f", None, "none"

# ---------- 분석 + 차트 + 설명 수집 ----------
def render_excel(file, file_name, prefs):
    df = pd.read_excel(file)
    st.subheader(f"📊 {file_name} 분석 결과")

    # 추천 플랜
    plan = recommend_plan(df)
    with st.expander("🔎 추천 보고서 플랜", expanded=True):
        st.markdown(_humanize_plan(plan))
        if use_ollama:
            one_liner = _ollama_generate(
                ollama_base, ollama_model,
                "다음 추천 플랜을 한국어 한 문장으로 요약해 주세요(존댓말): " + str(plan)
            )
            st.caption("🧠 자동 요약: " + one_liner)

    # 위젯 key 충돌 방지를 위한 prefix
    safe_name = file_name.replace("/", "_").replace("\\", "_")
    plan = _plan_editor(plan, key_prefix=f"plan_{safe_name}")

    # Ollama 개요
    overview_text = ""
    if use_ollama:
        prompt = (
            "다음 데이터의 컬럼과 추천 플랜을 바탕으로 한국어로 간결한 보고서 개요를 작성하세요.\n"
            f"- 대상: {prefs['audience']}\n- 어조: {prefs['tone']}\n"
            f"- 컬럼: {list(df.columns)}\n- 추천플랜: {plan}\n"
            "2~4문장으로 핵심 흐름(주요 지표, 추세, 비교)을 제안해 주세요."
        )
        overview_text = _ollama_generate(ollama_base, ollama_model, prompt)
        with st.expander("🗒️ 자동 개요", expanded=True):
            st.write(overview_text)

    chart_images = []
    chart_explanations = {}

    # 기본 통계
    st.write("✅ 데이터 요약")
    try:
        st.markdown(df.describe(include="all").to_markdown())
    except Exception:
        st.text(df.describe(include="all").to_string())

    # ---- 시계열 ----
    if prefs["sections"]["timeseries"] and plan["timeseries"]:
        for (dcol, ncol) in plan["timeseries"]:
            try:
                tdf = df[[dcol, ncol]].dropna()
                tdf[dcol] = pd.to_datetime(tdf[dcol], errors="coerce"); tdf = tdf.dropna().sort_values(dcol)
                title = f"{file_name}<br>{ncol} - 시계열({dcol})"
                fig = px.line(tdf, x=dcol, y=ncol, markers=True)
                _style_plotly(fig, title=title)
                st.plotly_chart(fig, use_container_width=True)

                exp = ""
                if use_ollama:
                    p = (f"시계열 그래프를 2~3문장으로 요약하세요. 대상:{prefs['audience']} 어조:{prefs['tone']} "
                         f"데이터셋:{file_name} x:{dcol} y:{ncol}. 주요 추세/피크/변동을 언급하세요.")
                    exp = _ollama_generate(ollama_base, ollama_model, p)
                    with st.expander(f"🗒️ {ncol} 시계열 설명", expanded=False):
                        st.write(exp)
                png = _fig_to_png_bytes(fig) or _mpl_hist_png(tdf[ncol], f"{file_name} · {ncol} 분포(대체)", ncol, 20)
                chart_images.append((f"{ncol} 시계열", png))
                chart_explanations[f"{ncol} 시계열"] = exp
            except Exception:
                pass

    # ---- 숫자형 분포 ----
    for col in plan["numeric_dists"]:
        histnorm = y_scale
        if show_kde and histnorm in ("count", None): histnorm = "probability density"
        title = f"{file_name}<br>{col} 분포"
        fig = px.histogram(df, x=col,
                           nbins=nbins if BIN_MODE=="개수 지정" else None,
                           color_discrete_sequence=["#4C78A8"],
                           histnorm=None if histnorm=="count" else histnorm)
        if BIN_MODE == "간격 지정" and binsize and binsize>0:
            fig.update_traces(xbins=dict(size=binsize))
        fig.update_traces(marker_line_color="white", marker_line_width=1, opacity=0.9)
        fig.update_layout(bargap=bargap, bargroupgap=0.06)
        y_label, y_tickformat, y_suffix, expfmt = _y_axis_label_and_format(histnorm)
        fig.update_xaxes(title_text=col)
        fig.update_yaxes(title_text=y_label, tickformat=y_tickformat,
                         ticksuffix=(y_suffix or ""), showexponent="none", exponentformat=expfmt)
        _style_plotly(fig, title=title)

        if show_kde:
            vals = df[col].dropna().values
            xs, ys = _try_kde(vals)
            if xs is not None and ys is not None:
                fig.add_scatter(x=xs, y=ys, mode="lines", name="KDE", line=dict(color="#E45756", width=2))
        st.plotly_chart(fig, use_container_width=True)

        exp = ""
        if use_ollama:
            note = " (핵심 KPI)" if col in prefs["kpis"] else ""
            prompt = (
                f"히스토그램 해설{note}: 데이터셋:{file_name}, 컬럼:{col}, 단위:{y_label}. "
                "치우침/범위/이상치를 2~3문장으로 간결히 설명하고, 업무 인사이트 1문장 제안(존댓말)."
            )
            exp = _ollama_generate(ollama_base, ollama_model, prompt)
            with st.expander(f"🗒️ {col} 분포 설명", expanded=False): st.write(exp)

        png = _fig_to_png_bytes(fig)
        if png is None:
            fallback_bins = nbins if (BIN_MODE=="개수 지정" and nbins) else 20
            png = _mpl_hist_png(df[col], f"{file_name} · {col} 분포", col, bins=fallback_bins)
        chart_images.append((f"{col} 분포", png))
        chart_explanations[f"{col} 분포"] = exp

    # ---- 범주형 분포 ----
    if "categoricals" in plan and plan["categoricals"]:
        for col in plan["categoricals"]:
            vc = df[col].astype(str).value_counts().head(15)
            title = f"{file_name}<br>{col} 상위 빈도"
            fig = px.bar(x=vc.index, y=vc.values, text=vc.values, labels={"x": col, "y": "개수"},
                         color_discrete_sequence=["#4C78A8"])
            _style_plotly(fig, title=title)
            fig.update_traces(textposition="outside"); fig.update_layout(yaxis_title="개수")
            st.plotly_chart(fig, use_container_width=True)

            exp = ""
            if use_ollama:
                prompt = (f"막대그래프 해설: 데이터셋:{file_name}, 범주:{col}, 상위 항목과 편중을 2문장으로 설명하고 "
                          f"업무적 시사점을 1문장 제안해 주세요(존댓말).")
                exp = _ollama_generate(ollama_base, ollama_model, prompt)
                with st.expander(f"🗒️ {col} 범주 설명", expanded=False): st.write(exp)

            png = _fig_to_png_bytes(fig) or _mpl_hist_png(pd.Series(vc.values), f"{col} 빈도(대체)", "빈도", 15)
            chart_images.append((f"{col} 범주", png))
            chart_explanations[f"{col} 범주"] = exp

    # ---- 상관관계 ----
    if prefs["sections"]["correlation"] and plan["correlation"]:
        num_cols = df.select_dtypes(include="number").columns
        corr = df[num_cols].corr(numeric_only=True)
        z = corr.values; x = corr.columns.tolist(); y = corr.columns.tolist()
        ann = corr.round(2).values
        heat = ff.create_annotated_heatmap(z=z, x=x, y=y, annotation_text=ann,
                                           colorscale="RdBu", showscale=True, reversescale=True)
        heat.update_coloraxes(colorbar_title="상관계수")
        _style_plotly(heat, title=f"{file_name}<br>숫자형 상관관계")
        st.plotly_chart(heat, use_container_width=True)

        exp = ""
        if use_ollama:
            pairs = []
            for i in range(len(x)):
                for j in range(i+1, len(y)):
                    pairs.append((x[i], y[j], float(corr.iloc[i,j])))
            pairs_sorted = sorted(pairs, key=lambda t: abs(t[2]), reverse=True)[:5]
            prompt = (f"상관 히트맵 해설: 데이터셋:{file_name}, 상위쌍:{pairs_sorted}. "
                      "유의미한 양/음의 상관과 주의점 2~3문장, 활용 제안 1문장을 한국어 존댓말로.")
            exp = _ollama_generate(ollama_base, ollama_model, prompt)
            with st.expander("🗒️ 상관관계 설명", expanded=False): st.write(exp)

        png = _fig_to_png_bytes(heat) or _mpl_corr_png(corr, f"{file_name} · 숫자형 상관관계")
        chart_images.append(("숫자형 상관관계", png))
        chart_explanations["숫자형 상관관계"] = exp

    # 사람 친화 플랜 텍스트 반환(보고서 개요에 포함 가능)
    plan_md_text = _humanize_plan(plan)
    return df, chart_images, chart_explanations, overview_text, plan_md_text

# ---------- 내보내기 ----------
def make_ppt_report(title: str, charts: dict, explanations: dict, overview_text: str) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"자동 생성 · {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    if overview_text:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "요약 개요"
        tx = s.shapes.add_textbox(PptxInches(0.8), PptxInches(1.5), PptxInches(8.4), PptxInches(3.6))
        tf = tx.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = overview_text; p.font.size = Pt(16)

    for dataset_name, items in charts.items():
        s = prs.slides.add_slide(prs.slide_layouts[5]); s.shapes.title.text = f"📦 {dataset_name}"
        for title, png_bytes in items:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
            if png_bytes:
                left = PptxInches(0.6); top = PptxInches(1.2); width = PptxInches(8.6)
                slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width=width)
            exp = explanations.get(dataset_name, {}).get(title, "")
            if exp:
                tx = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(5.7), PptxInches(8.6), PptxInches(1.6))
                tf = tx.text_frame; tf.clear()
                p = tf.paragraphs[0]; p.text = exp; p.font.size = Pt(14); p.alignment = PP_ALIGN.LEFT

    out = io.BytesIO(); prs.save(out); return out.getvalue()

def make_word_report(title: str, dfs: list[pd.DataFrame], texts: list[str],
                     charts: dict, explanations: dict, overview_text: str) -> bytes:
    doc = Document()
    doc.add_heading(title, level=1)
    doc.add_paragraph(f"자동 생성 · {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    if overview_text:
        doc.add_heading("요약 개요", level=2); doc.add_paragraph(overview_text)

    for i, df in enumerate(dfs, start=1):
        doc.add_heading(f"데이터셋 {i} 요약", level=2)
        desc = df.describe(include="all")
        table = doc.add_table(rows=1 + len(desc.index), cols=1 + len(desc.columns))
        table.style = "Light List Accent 1"
        hdr = table.rows[0].cells; hdr[0].text = ""
        for j, col in enumerate(desc.columns, start=1): hdr[j].text = str(col)
        for r_idx, idx_name in enumerate(desc.index, start=1):
            row = table.rows[r_idx].cells; row[0].text = str(idx_name)
            for c_idx, col in enumerate(desc.columns, start=1):
                val = desc.loc[idx_name, col]; row[c_idx].text = "" if pd.isna(val) else str(val)

    for ds_name, items in charts.items():
        doc.add_heading(f"차트 - {ds_name}", level=2)
        for (title, png_bytes) in items:
            if png_bytes:
                doc.add_paragraph(f"• {title}")
                doc.add_picture(io.BytesIO(png_bytes), width=DocxInches(6.5))
                exp = explanations.get(ds_name, {}).get(title, "")
                if exp: doc.add_paragraph(exp)

    for i, txt in enumerate(texts, start=1):
        doc.add_heading(f"PDF 문서 {i}", level=2)
        doc.add_paragraph(txt[:1500] + ("..." if len(txt) > 1500 else ""))

    out = io.BytesIO(); doc.save(out); return out.getvalue()

def make_excel_with_images(dfs, charts) -> bytes:
    out = io.BytesIO()
    with pd.ExcelWriter(out, engine="xlsxwriter") as writer:
        for i, df in enumerate(dfs, start=1):
            sheet_name = f"Data_{i}"
            df.to_excel(writer, sheet_name=sheet_name, index=False)
            ws = writer.sheets[sheet_name]
            for col_idx, col in enumerate(df.columns):
                max_len = max([len(str(col))] + [len(str(x)) for x in df[col].head(100).astype(str).tolist()])
                ws.set_column(col_idx, col_idx, min(max_len + 2, 40))
        chart_ws = writer.book.add_worksheet("Charts")
        for c in range(0, 24): chart_ws.set_column(c, c, 14)
        r, c, per_row, idx_in_row = 1, 1, 2, 0
        for ds_name, items in charts.items():
            chart_ws.write(r, c, f"📦 {ds_name}"); r += 1
            for (title, png_bytes) in items:
                if not png_bytes: continue
                chart_ws.write(r, c, f"• {title}")
                chart_ws.insert_image(r+1, c, "chart.png",
                                      {"image_data": io.BytesIO(png_bytes), "x_scale": 1.0, "y_scale": 1.0})
                idx_in_row += 1
                if idx_in_row % per_row == 0: r += 20; c = 1
                else: c += 8
            r += 22; c = 1; idx_in_row = 0
    return out.getvalue()

# ---------- PDF 텍스트 추출 ----------
def analyze_pdf(file, file_name):
    st.subheader(f"📄 {file_name} 텍스트 추출")
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text += (page.extract_text() or "") + "\n"
    st.text_area("📑 추출된 텍스트", text, height=180)
    return text

# ================== 메인 ==================
all_dfs, all_texts = [], []
all_charts: dict[str, list[tuple[str, bytes]]] = {}
all_explanations: dict[str, dict[str, str]] = {}
overview_by_dataset: dict[str, str] = {}
plan_text_by_dataset: dict[str, str] = {}

if uploaded_files:
    if len(st.session_state.messages) == 0:
        st.session_state.messages.append({"role":"assistant",
            "content":"어떤 스타일의 보고서를 원하세요? 예) '임원용, PPT/Word 중심, KPI는 매출·전환율, 상세'\n명령형: /audience 임원 /tone 상세 /out ppt,word /kpi 매출,전환율"})

    for m in st.session_state.messages:
        with st.chat_message(m["role"]):
            st.write(m["content"])
    if prompt := st.chat_input("보고서 요구사항을 입력하세요(명령어 지원: /audience, /tone, /out, /kpi)"):
        st.session_state.messages.append({"role":"user","content":prompt})
        bot = chat_bot_reply(prompt)
        st.session_state.messages.append({"role":"assistant","content":bot})
        with st.chat_message("assistant"):
            st.write(bot)

    for file in uploaded_files:
        name = file.name
        if name.lower().endswith(("xlsx","xls")):
            df, imgs, exps, ov, plan_md_text = render_excel(file, name, st.session_state.prefs)
            all_dfs.append(df)
            all_charts[name] = imgs
            all_explanations[name] = exps  # dict(title->exp)
            overview_by_dataset[name] = ov
            plan_text_by_dataset[name] = plan_md_text
        elif name.lower().endswith("pdf"):
            text = analyze_pdf(file, name)
            all_texts.append(text)

    if st.button("📥 결과 보고서 생성"):
        merged_overview = "\n\n".join([t for t in list(overview_by_dataset.values()) + list(plan_text_by_dataset.values()) if t])

        if st.session_state.prefs["outputs"]["md"]:
            md = "# 🎯 이벤트 결과 보고서\n\n"
            md += f"- 대상: {st.session_state.prefs['audience']} / 어조: {st.session_state.prefs['tone']}\n"
            md += "✨ 자동 생성된 요약 리포트입니다.\n\n"
            if merged_overview:
                md += "## 개요\n" + merged_overview + "\n\n"
            for i, df in enumerate(all_dfs, start=1):
                md += f"## 데이터셋 {i} 요약\n"
                try: md += df.describe(include="all").to_markdown() + "\n\n"
                except Exception: md += df.describe(include="all").to_string() + "\n\n"
            st.download_button("📥 Markdown 보고서 다운로드", data=md, file_name="event_report.md")

        if st.session_state.prefs["outputs"]["excel"]:
            xlsx = make_excel_with_images(all_dfs, all_charts)
            st.download_button("📊 Excel 보고서(차트 내장) 다운로드", data=xlsx, file_name="event_report_with_charts.xlsx")

        if st.session_state.prefs["outputs"]["ppt"]:
            ppt = make_ppt_report("이벤트 결과 보고서 ✨", all_charts, all_explanations, merged_overview)
            st.download_button("📽 PPT 보고서(차트·설명 포함) 다운로드", data=ppt, file_name="event_report.pptx")

        if st.session_state.prefs["outputs"]["docx"]:
            docx = make_word_report("이벤트 결과 보고서 ✨", all_dfs, all_texts, all_charts, all_explanations, merged_overview)
            st.download_button("📝 Word 보고서(.docx) 다운로드", data=docx, file_name="event_report.docx")

        if any(len(v)>0 for v in all_charts.values()):
            zip_buf = io.BytesIO()
            with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
                for ds, items in all_charts.items():
                    safe = ds.replace("/", "_")
                    for i, (ctitle, png) in enumerate(items, start=1):
                        if not png: continue
                        zf.writestr(f"{safe}/chart_{i:02d}_{ctitle}.png", png)
            st.download_button("🖼 차트 PNG 묶음(ZIP) 다운로드", data=zip_buf.getvalue(), file_name="charts.zip")
